#!/usr/bin/env python3
"""
TRIO 3.0 · PPTX CJK 字体后处理
用法: python3 pptx-cjk-lock.py <input.pptx> [output.pptx]
功能: 遍历所有幻灯片，将所有文本的字体锁定为 Noto Sans CJK SC
"""
import sys
from pptx import Presentation
from pptx.oxml.ns import qn

FONT_FAMILY = "Noto Sans CJK SC"
MONO_FAMILY = "Noto Sans Mono"


def lock_fonts(pptx_path, output_path=None):
    """锁定 PPTX 中所有文本的 CJK 字体"""
    prs = Presentation(pptx_path)
    fixed = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    # 设置拉丁字体
                    run.font.name = FONT_FAMILY
                    # 设置东亚字体（关键——pandoc 不会设这个）
                    rPr = run._r.get_or_add_rPr()
                    # 移除旧的东西字体设置
                    for old in rPr.findall(qn('a:ea')):
                        rPr.remove(old)
                    ea = rPr.makeelement(qn('a:ea'), {})
                    ea.set('typeface', FONT_FAMILY)
                    rPr.append(ea)
                    fixed += 1

    out = output_path or pptx_path
    prs.save(out)
    print(f"✅ CJK 字体已锁定 → {FONT_FAMILY}")
    print(f"   处理: {len(prs.slides)} 页 | {fixed} 个文本运行")
    print(f"   输出: {out}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python3 pptx-cjk-lock.py <input.pptx> [output.pptx]")
        sys.exit(1)
    lock_fonts(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
